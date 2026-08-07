import os
import json
import signal
from pathlib import Path
import pandas as pd
from pptx import Presentation
from docx import Document as DocxDocument
from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document


class TimeoutError(Exception):
    pass

def _timeout_handler(signum, frame):
    raise TimeoutError()

def load_pdf(file_path):
    return PyPDFLoader(file_path).load()

def load_txt_or_code(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return [Document(page_content=text, metadata={"source": file_path})]

def load_pptx(file_path):
    prs = Presentation(file_path)
    docs = []
    for i, slide in enumerate(prs.slides):
        text = "\n".join(
            shape.text for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
        if text.strip():
            docs.append(Document(
                page_content=f"[Slide {i+1}]\n{text}",
                metadata={"source": file_path, "slide": i+1}
            ))
    return docs

def load_docx(file_path):
    doc = DocxDocument(file_path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [Document(page_content=text, metadata={"source": file_path})]

def load_csv(file_path):
    df = pd.read_csv(file_path, encoding="utf-8", errors="ignore")
    return [Document(page_content=df.to_string(index=False), metadata={"source": file_path})]

def load_xlsx(file_path):
    xl = pd.ExcelFile(file_path)
    all_text = []
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        all_text.append(f"[Sheet: {sheet}]\n{df.to_string(index=False)}")
    return [Document(page_content="\n\n".join(all_text), metadata={"source": file_path})]

def load_json(file_path):
    try:
        data = json.loads(Path(file_path).read_text())
        if "history" not in data:
            return []
        docs = []
        for turn in data["history"]:
            text = (
                f"[Past conversation — {data.get('started', '')}]\n"
                f"Q: {turn['question']}\nA: {turn['answer']}"
            )
            docs.append(Document(
                page_content=text,
                metadata={"source": str(file_path), "type": "conversation"}
            ))
        return docs
    except:
        return []

def load_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    dispatch = {
        ".pdf":  load_pdf,
        ".pptx": load_pptx,
        ".docx": load_docx,
        ".csv":  load_csv,
        ".xlsx": load_xlsx,
        ".json": load_json,
    }
    code_exts = {".txt", ".py", ".js", ".ts", ".cpp", ".c", ".h", ".java", ".md"}
    if ext in dispatch:
        return dispatch[ext](file_path)
    elif ext in code_exts:
        return load_txt_or_code(file_path)
    return []

def load_file_with_timeout(file_path, seconds=60):
    signal.signal(signal.SIGALRM, _timeout_handler)
    signal.alarm(seconds)
    try:
        docs = load_file(file_path)
        signal.alarm(0)
        return docs
    except TimeoutError:
        print(f"  ⏱ Skipped (>{seconds}s): {os.path.basename(file_path)}")
    except Exception as e:
        print(f"  ✗ Error: {os.path.basename(file_path)} — {e}")
    # Return dummy so it's never retried
    return [Document(
        page_content="[SKIPPED — failed during indexing]",
        metadata={"source": file_path, "skipped": True}
    )]
